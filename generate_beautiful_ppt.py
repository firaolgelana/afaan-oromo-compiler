from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

output_path = r"C:\Users\sebon\OneDrive\Desktop\Afaan_Oromoo_Compiler_Presentation.pptx"

prs = Presentation()

# --- Design Colors ---
BG_COLOR = RGBColor(30, 30, 30)          # Dark grey/black
ACCENT_COLOR = RGBColor(0, 150, 255)     # Bright Blue
TITLE_COLOR = RGBColor(255, 204, 0)      # Gold/Yellow
TEXT_COLOR = RGBColor(240, 240, 240)     # Off-white

def apply_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header_accent(slide):
    # Add a thin blue line under the title
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(0.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()

def format_title(title_shape, text, is_main=False):
    title_shape.text = text
    if not title_shape.text_frame.paragraphs:
        return
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = 'Segoe UI'
    if is_main:
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
    else:
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR

def format_body(body_shape, content_list):
    tf = body_shape.text_frame
    tf.clear()  # Clear existing
    for i, content in enumerate(content_list):
        p = tf.add_paragraph()
        if content.startswith("  -"):
            p.text = content.replace("  - ", "").strip()
            p.level = 1
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_COLOR
        elif content.startswith("    "):
            p.text = content.strip()
            p.level = 2
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(200, 200, 200) # Slightly dimmer
            p.font.name = 'Consolas' # Code font
        elif content.startswith("caasaa") or content.startswith("gocha") or content.startswith("yoo ") or content.startswith("}") or content.startswith("hanga ") or content.startswith("marsaa "):
            p.text = content.strip()
            p.level = 1
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(180, 255, 180) # Code color
            p.font.name = 'Consolas'
        else:
            p.text = content.strip()
            p.level = 0
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR
        p.font.name = 'Segoe UI' if p.font.name != 'Consolas' else 'Consolas'

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide)
    
    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    format_title(txBox, title, is_main=True)
    
    # Add subtitle
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox2.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    
    add_header_accent(slide)

def add_content_slide(title, content_list):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide)
    add_header_accent(slide)
    
    # Add title
    left = Inches(1.0)
    top = Inches(0.5)
    width = Inches(8.0)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    format_title(txBox, title, is_main=False)
    
    # Add body
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(5.5)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    txBox2.text_frame.word_wrap = True
    format_body(txBox2, content_list)

# 1
add_title_slide(
    "Afaan Oromoo Compiler",
    "Building a Native Programming Language\nusing a Source-to-Source Architecture"
)

# 2
add_content_slide(
    "1. Introduction",
    [
        "What is the Afaan Oromoo Compiler?",
        "  - A fully functional programming language written in and for Afaan Oromoo.",
        "  - Designed to bridge the language gap in computer science education.",
        "  - Allows developers to write code using their native language keywords and syntax.",
        "  - Compiles seamlessly to executable Python code."
    ]
)

# 3
add_content_slide(
    "2. Motivation & Goals",
    [
        "Why create an Afaan Oromoo programming language?",
        "  - Accessibility: Lower the barrier to entry for native speakers learning to code.",
        "  - Expressiveness: Use familiar terminology for logic and structures.",
        "  - Performance: Leverage existing highly-optimized ecosystems.",
        "  - Education: Teach core computer science concepts without the English barrier."
    ]
)

# 4
add_content_slide(
    "3. Core Architecture: The Transpiler",
    [
        "What is a Source-to-Source Compiler (Transpiler)?",
        "  - Translates source code from one high-level programming language to another.",
        "  - Afaan Oromoo -> Python.",
        "  - It does NOT compile directly to machine code or bytecode.",
        "  - This allows us to stand on the shoulders of giants (Python's runtime)."
    ]
)

# 5
add_content_slide(
    "4. Why Transpile to Python?",
    [
        "Strategic Advantages:",
        "  - Ecosystem: Instant access to thousands of Python libraries (AI, Data, Web).",
        "  - Cross-Platform: Runs anywhere Python runs (Windows, Mac, Linux).",
        "  - Execution Speed: Python's heavily optimized C-backend.",
        "  - Simplicity: Allows us to focus on Afaan Oromoo syntax and semantics."
    ]
)

# 6
add_content_slide(
    "5. The Compilation Pipeline",
    [
        "How Afaan Oromoo code becomes executable:",
        "  - 1. Lexical Analysis (Scanner): Converts raw text into Tokens.",
        "  - 2. Parsing: Builds an Abstract Syntax Tree (AST).",
        "  - 3. Code Generation (Codegen): Traverses AST to emit Python 3 source code.",
        "  - 4. Execution: Runs the generated Python code natively."
    ]
)

# 7
add_content_slide(
    "6. Syntax Style Overview",
    [
        "A Unique Blend of Python and C++",
        "  - Python's Influence: Clean, readable, dynamically typed, rich data structures.",
        "  - C++'s Influence: Structured block scopes using curly braces { }.",
        "  - The Extra Mile: Native Afaan Oromoo keywords and culturally relevant semantics."
    ]
)

# 8
add_content_slide(
    "7. Syntax Style: Python's Influence",
    [
        "What we borrowed from Python:",
        "  - Dynamic Typing: Variables do not need explicit types.",
        "  - High-level Data Structures: Native Arrays (tarree) and Dictionaries.",
        "  - Control Flow Semantics: If, Elif, Else behave exactly like Python.",
        "  - Iteration: For-in loops over iterables."
    ]
)

# 9
add_content_slide(
    "8. Syntax Style: C++'s Influence",
    [
        "What we borrowed from C/C++:",
        "  - Explicit Blocks: We use curly braces '{' and '}' instead of whitespace indentation.",
        "  - Why? Curly braces are often more explicit for beginners to see where a block begins and ends.",
        "yoo (x > 5) {",
        "    maxxansi(x)",
        "}",
        "  - Increment/Decrement: Native support for ++ and -- operators."
    ]
)

# 10
add_content_slide(
    "9. Lexical Analysis (Tokens)",
    [
        "Translating text into meaning:",
        "  - Keywords: yoo (if), yookaas (elif), yookiin (else), hanga (while), marsaa (for).",
        "  - Operators: +, -, *, /, %, **, //, ==, !=, <=, >=",
        "  - Delimiters: {, }, (, ), [, ]",
        "  - Identifiers: Variable names and function names."
    ]
)

# 11
add_content_slide(
    "10. Parsing (Abstract Syntax Tree)",
    [
        "Structuring the grammar:",
        "  - Uses Recursive Descent Parsing.",
        "  - Nodes include: BinaryOp, If, ForLoop, ClassDef, ListComprehension.",
        "  - Validates syntax rules and structural integrity before any code is generated."
    ]
)

# 12
add_content_slide(
    "11. Code Generation",
    [
        "Translating AST to Python:",
        "  - Walks the AST using the Visitor Pattern.",
        "  - Emits properly indented Python 3 code.",
        "  - Automatically translates built-in functions (e.g. maxxansi -> print).",
        "  - Automatically translates object methods (e.g. dabali -> append)."
    ]
)

# 13
add_content_slide(
    "12. Data Types: Fundamentals",
    [
        "Native support for basic types:",
        "  - Numbers: Integers (10) and Floats (3.14).",
        "  - Strings: Single and Double quoted text (\"Afaan Oromoo\").",
        "  - Multi-line Strings: Supported via triple quotes \"\"\" ... \"\"\".",
        "  - Operators: Full math suite including ** (power) and // (floor division)."
    ]
)

# 14
add_content_slide(
    "13. Data Types: Logic",
    [
        "Boolean and Null values localized:",
        "  - True -> dhugaa",
        "  - False -> soba",
        "  - None (Null) -> homaa",
        "yoo (dhugaa) {",
        "    maxxansi(\"Sirrii!\")",
        "}"
    ]
)

# 15
add_content_slide(
    "14. Data Structures: Collections",
    [
        "Complex data organization:",
        "  - Arrays (tarree): [1, 2, 3]",
        "  - Dictionaries: {\"maqaa\": \"Abdii\", \"umrii\": 25}",
        "  - Tuples (tuuta_dhaabbataa): (10, 20)",
        "  - Sets (tuuta): {1, 2, 3}"
    ]
)

# 16
add_content_slide(
    "15. Control Flow: Branching",
    [
        "Making decisions in code:",
        "yoo (x > 10) {",
        "    maxxansi(\"Guddaa dha\")",
        "} yookaas (x == 10) {",
        "    maxxansi(\"Qixa\")",
        "} yookiin {",
        "    maxxansi(\"Xiqqaa dha\")",
        "}"
    ]
)

# 17
add_content_slide(
    "16. Control Flow: Looping",
    [
        "Repeating actions:",
        "While loop:",
        "hanga (x < 5) {",
        "    x++",
        "}",
        "For loop (Iteration):",
        "marsaa (x keessa tarree) {",
        "    maxxansi(x)",
        "}"
    ]
)

# 18
add_content_slide(
    "17. Functions",
    [
        "Defining reusable logic:",
        "  - Keyword: gocha (or hojii)",
        "gocha ida'uu(a, b) {",
        "    deebisi a + b",
        "}"
    ]
)

# 19
add_content_slide(
    "18. Object-Oriented Programming",
    [
        "Creating custom structures:",
        "  - Keyword: caasaa (class)",
        "  - Supports encapsulation and behavior.",
        "caasaa Nama {",
        "    gocha dubbadhu(ofii) {",
        "        maxxansi(\"Akkam!\")",
        "    }",
        "}"
    ]
)

# 20
add_content_slide(
    "19. OOP: Inheritance & Constructors",
    [
        "Advanced Class Features:",
        "  - Constructor: Uses the __jalqaba__ method (maps to __init__).",
        "  - Instance reference: Uses 'ofii' (maps to self).",
        "  - Inheritance:",
        "caasaa Barataa(Nama) { ... }"
    ]
)

# 21
add_content_slide(
    "20. Core Built-in Functions",
    [
        "Essential functions provided globally:",
        "  - I/O: maxxansi (print), galchi (input)",
        "  - Type conversion: lakkoofsa (int), barruu (str), desimaalii (float)",
        "  - Math: guddaa (max), xiqqaa (min), naannessi (round), gatsirrii (abs)",
        "  - File System: bani (open)"
    ]
)

# 22
add_content_slide(
    "21. Method Operations",
    [
        "Native method calls using dot notation:",
        "  - Array ops: dabali (append), baasi (pop), haqi (remove), tartiibessi (sort)",
        "  - String ops: qubeeguddaa (upper), qoqqoodi (split)",
        "  - Dict ops: furtuuwwan (keys), gatiiwwan (values), argadhu (get)",
        "  - The transpiler dynamically maps these during codegen."
    ]
)

# 23
add_content_slide(
    "22. Advanced Feature: List Comprehensions",
    [
        "Elegant array transformations:",
        "  - Python style: [x * 2 for x in nums if x > 0]",
        "  - Afaan Oromoo style:",
        "    [x * 2 marsaa x keessa nums yoo x > 0]",
        "  - High-performance, declarative data processing."
    ]
)

# 24
add_content_slide(
    "23. Advanced Helpers",
    [
        "Global helper functions for data manipulation:",
        "  - qindeessi: returns a sorted version of an iterable.",
        "  - lakkoofsa_waliin: (enumerate) adds a counter to an iterable.",
        "  - walitti_hidhi: (zip) aggregates elements from multiple iterables."
    ]
)

# 25
add_content_slide(
    "24. Modules and Imports",
    [
        "Extending the language:",
        "  - Since we transpile to Python, we can import Python modules directly!",
        "  - Mathematical libraries, AI tools, web frameworks.",
        "  - The 'import' system resolves module names and handles native bridging automatically."
    ]
)

# 26
add_content_slide(
    "25. Import Syntax",
    [
        "How to load modules:",
        "  - Standard import: fidi math",
        "  - Import all: irraa math fidi *",
        "  - Allows Afaan Oromoo developers to use packages like NumPy or TensorFlow directly."
    ]
)

# 27
add_content_slide(
    "26. IDE Integration: VS Code Extension",
    [
        "A language needs proper tooling:",
        "  - We developed a dedicated VS Code Extension (vscode-oromolang).",
        "  - Provides rich syntax highlighting for keywords, operators, and built-ins.",
        "  - Uses TextMate grammar JSON format.",
        "  - Dramatically improves developer experience."
    ]
)

# 28
add_content_slide(
    "27. Error Handling & Diagnostics",
    [
        "Debugging made easier:",
        "  - Custom Error Messages localized in Afaan Oromoo.",
        "  - E.g., 'Dogoggora: Mallattoo hin eegamne' (Unexpected token).",
        "  - Shows exact line and column (giddugaleessa) where syntax errors occur."
    ]
)

# 29
add_content_slide(
    "28. Testing and Validation",
    [
        "Ensuring language stability:",
        "  - Multi-phase implementation (P1, P2, P3).",
        "  - Comprehensive test suites (test_p1.ao, test_p2.ao, test_p3.ao).",
        "  - Verifies AST construction and proper Python code generation."
    ]
)

# 30
add_content_slide(
    "29. Future Work & Roadmap",
    [
        "What's next for the language?",
        "  - Standard Library: Building wrappers for Python libraries purely in Afaan Oromoo.",
        "  - F-Strings: Implementing native string interpolation.",
        "  - Community: Open sourcing to the Ethiopian tech community.",
        "  - Package Manager: A custom tool for downloading Afaan Oromoo packages."
    ]
)

# 31
add_content_slide(
    "30. Conclusion",
    [
        "Summary:",
        "  - Built a powerful, culturally relevant language.",
        "  - Used Source-to-Source transpilation to leverage Python.",
        "  - Combined Python's expressiveness with C++ structural blocks.",
        "",
        "Galatoomaa! (Thank you!)",
        "Questions?"
    ]
)

prs.save(output_path)
print("Beautiful presentation generated successfully!")
