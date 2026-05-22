from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

template_path = r"C:\Users\sebon\OneDrive\Desktop\SmartPharmacy_Presentation.pptx"
output_path = r"C:\Users\sebon\OneDrive\Desktop\Afaan_Oromoo_Compiler_Presentation.pptx"

prs = Presentation(template_path)

# Delete existing slides
xml_slides = prs.slides._sldIdLst  
slides = list(xml_slides)
for slide in slides:
    xml_slides.remove(slide)

def add_title_slide(title, subtitle):
    # Try to use layout 0, fallback to 0 if not enough layouts
    layout_index = 0
    if len(prs.slide_layouts) == 0:
        return
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Try to set subtitle
    if len(slide.placeholders) > 1:
        # Find a placeholder that isn't the title
        for shape in slide.placeholders:
            if shape != slide.shapes.title:
                shape.text = subtitle
                break

def add_content_slide(title, content_list):
    layout_index = 1 if len(prs.slide_layouts) > 1 else 0
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
        
    body_shape = None
    for shape in slide.placeholders:
        if shape != slide.shapes.title:
            body_shape = shape
            break
            
    if body_shape and hasattr(body_shape, "text_frame"):
        tf = body_shape.text_frame
        for i, content in enumerate(content_list):
            if i == 0:
                tf.text = content
            else:
                p = tf.add_paragraph()
                if content.startswith("  "):
                    p.text = content.strip()
                    p.level = 1
                else:
                    p.text = content
                    p.level = 0

# 1
add_title_slide(
    "Afaan Oromoo Compiler",
    "Building a Native Programming Language\nusing a Source-to-Source Architecture"
)

# 2
add_content_slide(
    "1. Introduction",
    [
        "What is the Afaan Oromoo Compiler?",
        "  - A fully functional programming language written in and for Afaan Oromoo.",
        "  - Designed to bridge the language gap in computer science education.",
        "  - Allows developers to write code using their native language keywords and syntax.",
        "  - Compiles seamlessly to executable Python code."
    ]
)

# 3
add_content_slide(
    "2. Motivation & Goals",
    [
        "Why create an Afaan Oromoo programming language?",
        "  - Accessibility: Lower the barrier to entry for native speakers learning to code.",
        "  - Expressiveness: Use familiar terminology for logic and structures.",
        "  - Performance: Leverage existing highly-optimized ecosystems.",
        "  - Education: Teach core computer science concepts without the English barrier."
    ]
)

# 4
add_content_slide(
    "3. Core Architecture: The Transpiler",
    [
        "What is a Source-to-Source Compiler (Transpiler)?",
        "  - Translates source code from one high-level programming language to another.",
        "  - Afaan Oromoo -> Python.",
        "  - It does NOT compile directly to machine code or bytecode.",
        "  - This allows us to stand on the shoulders of giants (Python's runtime)."
    ]
)

# 5
add_content_slide(
    "4. Why Transpile to Python?",
    [
        "Strategic Advantages:",
        "  - Ecosystem: Instant access to thousands of Python libraries (AI, Data, Web).",
        "  - Cross-Platform: Runs anywhere Python runs (Windows, Mac, Linux).",
        "  - Execution Speed: Python's heavily optimized C-backend.",
        "  - Simplicity: Allows us to focus on Afaan Oromoo syntax and semantics."
    ]
)

# 6
add_content_slide(
    "5. The Compilation Pipeline",
    [
        "How Afaan Oromoo code becomes executable:",
        "  1. Lexical Analysis (Scanner): Converts raw text into Tokens.",
        "  2. Parsing: Builds an Abstract Syntax Tree (AST) representing the program's structure.",
        "  3. Code Generation (Codegen): Traverses the AST to emit Python 3 source code.",
        "  4. Execution: Runs the generated Python code natively."
    ]
)

# 7
add_content_slide(
    "6. Syntax Style Overview",
    [
        "A Unique Blend of Python and C++",
        "  - Python's Influence: Clean, readable, dynamically typed, rich data structures.",
        "  - C++'s Influence: Structured block scopes using curly braces { }.",
        "  - The Extra Mile: Native Afaan Oromoo keywords and culturally relevant semantics."
    ]
)

# 8
add_content_slide(
    "7. Syntax Style: Python's Influence",
    [
        "What we borrowed from Python:",
        "  - Dynamic Typing: Variables do not need explicit types.",
        "  - High-level Data Structures: Native Arrays (tarree) and Dictionaries.",
        "  - Control Flow Semantics: If, Elif, Else behave exactly like Python.",
        "  - Iteration: For-in loops over iterables."
    ]
)

# 9
add_content_slide(
    "8. Syntax Style: C++'s Influence",
    [
        "What we borrowed from C/C++:",
        "  - Explicit Blocks: We use curly braces '{' and '}' instead of whitespace indentation.",
        "  - Why? Curly braces are often more explicit for beginners to see where a block begins and ends.",
        "  - Example: yoo (x > 5) { maxxansi(x) }",
        "  - Increment/Decrement: Native support for ++ and -- operators."
    ]
)

# 10
add_content_slide(
    "9. Lexical Analysis (Tokens)",
    [
        "Translating text into meaning:",
        "  - Keywords: yoo (if), yookaas (elif), yookiin (else), hanga (while), marsaa (for).",
        "  - Operators: +, -, *, /, %, **, //, ==, !=, <=, >=",
        "  - Delimiters: {, }, (, ), [, ]",
        "  - Identifiers: Variable names and function names."
    ]
)

# 11
add_content_slide(
    "10. Parsing (Abstract Syntax Tree)",
    [
        "Structuring the grammar:",
        "  - Uses Recursive Descent Parsing.",
        "  - Nodes include: BinaryOp, If, ForLoop, ClassDef, ListComprehension.",
        "  - Validates syntax rules and structural integrity before any code is generated."
    ]
)

# 12
add_content_slide(
    "11. Code Generation",
    [
        "Translating AST to Python:",
        "  - Walks the AST using the Visitor Pattern.",
        "  - Emits properly indented Python 3 code.",
        "  - Automatically translates built-in functions (e.g. maxxansi -> print).",
        "  - Automatically translates object methods (e.g. dabali -> append)."
    ]
)

# 13
add_content_slide(
    "12. Data Types: Fundamentals",
    [
        "Native support for basic types:",
        "  - Numbers: Integers (10) and Floats (3.14).",
        "  - Strings: Single and Double quoted text (\"Afaan Oromoo\").",
        "  - Multi-line Strings: Supported via triple quotes \"\"\" ... \"\"\".",
        "  - Operators: Full math suite including ** (power) and // (floor division)."
    ]
)

# 14
add_content_slide(
    "13. Data Types: Logic",
    [
        "Boolean and Null values localized:",
        "  - True -> dhugaa",
        "  - False -> soba",
        "  - None (Null) -> homaa",
        "  - Example: yoo (dhugaa) { maxxansi(\"Sirrii!\") }"
    ]
)

# 15
add_content_slide(
    "14. Data Structures: Collections",
    [
        "Complex data organization:",
        "  - Arrays (tarree): [1, 2, 3]",
        "  - Dictionaries: {\"maqaa\": \"Abdii\", \"umrii\": 25}",
        "  - Tuples (tuuta_dhaabbataa): (10, 20)",
        "  - Sets (tuuta): {1, 2, 3}"
    ]
)

# 16
add_content_slide(
    "15. Control Flow: Branching",
    [
        "Making decisions in code:",
        "yoo (x > 10) {",
        "    maxxansi(\"Guddaa dha\")",
        "} yookaas (x == 10) {",
        "    maxxansi(\"Qixa\")",
        "} yookiin {",
        "    maxxansi(\"Xiqqaa dha\")",
        "}"
    ]
)

# 17
add_content_slide(
    "16. Control Flow: Looping",
    [
        "Repeating actions:",
        "While loop:",
        "hanga (x < 5) { x++ }",
        "",
        "For loop (Iteration):",
        "marsaa (x keessa tarree) { maxxansi(x) }"
    ]
)

# 18
add_content_slide(
    "17. Functions",
    [
        "Defining reusable logic:",
        "  - Keyword: gocha (or hojii)",
        "  - Example:",
        "gocha ida'uu(a, b) {",
        "    deebisi a + b",
        "}"
    ]
)

# 19
add_content_slide(
    "18. Object-Oriented Programming",
    [
        "Creating custom structures:",
        "  - Keyword: caasaa (class)",
        "  - Supports encapsulation and behavior.",
        "caasaa Nama {",
        "    gocha dubbadhu(ofii) {",
        "        maxxansi(\"Akkam!\")",
        "    }",
        "}"
    ]
)

# 20
add_content_slide(
    "19. OOP: Inheritance & Constructors",
    [
        "Advanced Class Features:",
        "  - Constructor: Uses the __jalqaba__ method (maps to __init__).",
        "  - Instance reference: Uses 'ofii' (maps to self).",
        "  - Inheritance: caasaa Barataa(Nama) { ... }",
        "  - Allows for highly structured, modern application design."
    ]
)

# 21
add_content_slide(
    "20. Core Built-in Functions",
    [
        "Essential functions provided globally:",
        "  - I/O: maxxansi (print), galchi (input)",
        "  - Type conversion: lakkoofsa (int), barruu (str), desimaalii (float)",
        "  - Math: guddaa (max), xiqqaa (min), naannessi (round), gatsirrii (abs)",
        "  - File System: bani (open)"
    ]
)

# 22
add_content_slide(
    "21. Method Operations",
    [
        "Native method calls using dot notation:",
        "  - Array ops: dabali (append), baasi (pop), haqi (remove), tartiibessi (sort)",
        "  - String ops: qubeeguddaa (upper), qoqqoodi (split)",
        "  - Dict ops: furtuuwwan (keys), gatiiwwan (values), argadhu (get)",
        "  - The transpiler dynamically maps these during codegen."
    ]
)

# 23
add_content_slide(
    "22. Advanced Feature: List Comprehensions",
    [
        "Elegant array transformations:",
        "  - Python style: [x * 2 for x in nums if x > 0]",
        "  - Afaan Oromoo style:",
        "    [x * 2 marsaa x keessa nums yoo x > 0]",
        "  - High-performance, declarative data processing."
    ]
)

# 24
add_content_slide(
    "23. Advanced Helpers",
    [
        "Global helper functions for data manipulation:",
        "  - qindeessi: returns a sorted version of an iterable.",
        "  - lakkoofsa_waliin: (enumerate) adds a counter to an iterable.",
        "  - walitti_hidhi: (zip) aggregates elements from multiple iterables."
    ]
)

# 25
add_content_slide(
    "24. Modules and Imports",
    [
        "Extending the language:",
        "  - Since we transpile to Python, we can import Python modules directly!",
        "  - Mathematical libraries, AI tools, web frameworks.",
        "  - The 'import' system resolves module names and handles native bridging automatically."
    ]
)

# 26
add_content_slide(
    "25. Import Syntax",
    [
        "How to load modules:",
        "  - Standard import: fidi math",
        "  - Import all: irraa math fidi *",
        "  - Allows Afaan Oromoo developers to use packages like NumPy or TensorFlow directly."
    ]
)

# 27
add_content_slide(
    "26. IDE Integration: VS Code Extension",
    [
        "A language needs proper tooling:",
        "  - We developed a dedicated VS Code Extension (vscode-oromolang).",
        "  - Provides rich syntax highlighting for keywords, operators, and built-ins.",
        "  - Uses TextMate grammar JSON format.",
        "  - Dramatically improves developer experience."
    ]
)

# 28
add_content_slide(
    "27. Error Handling & Diagnostics",
    [
        "Debugging made easier:",
        "  - Custom Error Messages localized in Afaan Oromoo.",
        "  - E.g., 'Dogoggora: Mallattoo hin eegamne' (Unexpected token).",
        "  - Shows exact line and column (giddugaleessa) where syntax errors occur."
    ]
)

# 29
add_content_slide(
    "28. Testing and Validation",
    [
        "Ensuring language stability:",
        "  - Multi-phase implementation (P1, P2, P3).",
        "  - Comprehensive test suites (test_p1.ao, test_p2.ao, test_p3.ao).",
        "  - Verifies AST construction and proper Python code generation."
    ]
)

# 30
add_content_slide(
    "29. Future Work & Roadmap",
    [
        "What's next for the language?",
        "  - Standard Library: Building wrappers for Python libraries purely in Afaan Oromoo.",
        "  - F-Strings: Implementing native string interpolation.",
        "  - Community: Open sourcing to the Ethiopian tech community.",
        "  - Package Manager: A custom tool for downloading Afaan Oromoo packages."
    ]
)

# 31
add_content_slide(
    "30. Conclusion",
    [
        "Summary:",
        "  - Built a powerful, culturally relevant language.",
        "  - Used Source-to-Source transpilation to leverage Python.",
        "  - Combined Python's expressiveness with C++ structural blocks.",
        "",
        "Galatoomaa! (Thank you!)",
        "Questions?"
    ]
)

prs.save(output_path)
print("Styled presentation generated successfully!")
